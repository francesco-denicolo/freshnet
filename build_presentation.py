# -*- coding: utf-8 -*-
"""
Genera una presentazione corporate (esteso, taglio misto) del paper
"Imputer x Forecaster Benchmark on FreshRetailNet-50K".
Fonte numeri: PAPER_FINDINGS.md (2026-06-12, 113 celle, Friedman + Kendall W + Nemenyi CD).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

FIG = "pipeline/figures"

# ---------------------------------------------------------------- palette
NAVY   = RGBColor(0x0E, 0x1F, 0x3A)   # sfondo divider / titoli
BLUE   = RGBColor(0x1F, 0x4E, 0x79)   # accent primario
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)   # accent secondario
AMBER  = RGBColor(0xE9, 0xA5, 0x2B)   # highlight
GREY   = RGBColor(0x5A, 0x5A, 0x5A)
LGREY  = RGBColor(0xEC, 0xEF, 0xF3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x20, 0x24, 0x2B)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)

FONT = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text,size,bold,color,italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for (t, sz, b, col, *rest) in para:
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = col; r.font.name = FONT
            r.font.italic = it
    return tb


def header(s, kicker, title, idx):
    """Standard content-slide header band."""
    rect(s, 0, 0, SW, Inches(1.15), WHITE)
    rect(s, 0, Inches(1.15), SW, Pt(3), BLUE)
    rect(s, Inches(0.55), Inches(0.30), Pt(5), Inches(0.62), AMBER)
    txt(s, Inches(0.75), Inches(0.22), Inches(11.4), Inches(0.32),
        [(kicker.upper(), 12, True, TEAL)])
    txt(s, Inches(0.75), Inches(0.50), Inches(11.6), Inches(0.55),
        [(title, 25, True, NAVY)])
    # page number
    txt(s, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
        [(str(idx), 10, False, GREY)], align=PP_ALIGN.RIGHT)
    rect(s, Inches(0.55), Inches(7.12), Inches(3.2), Pt(1.2), LGREY)
    txt(s, Inches(0.55), Inches(7.02), Inches(4.5), Inches(0.3),
        [("FreshRetailNet-50K  ·  Imputer × Forecaster Benchmark", 8, False, GREY)])


def bullets(s, x, y, w, h, items, size=15, gap=7):
    """items: list of (text, level, color_opt)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        if len(it) == 2:
            t, lvl = it; col = DARK; bold = False
        else:
            t, lvl, col = it[0], it[1], it[2]; bold = (lvl == 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.02
        mark = "▸ " if lvl == 0 else ("•  " if lvl == 1 else "–  ")
        p.level = 0
        run = p.add_run(); run.text = mark
        run.font.size = Pt(size); run.font.bold = (lvl == 0)
        run.font.color.rgb = (BLUE if lvl == 0 else TEAL)
        run.font.name = FONT
        # indentation via spaces for sub-levels
        prefix = "   " * lvl
        r2 = p.add_run(); r2.text = prefix + t
        r2.font.size = Pt(size - (0 if lvl == 0 else 1))
        r2.font.bold = bold
        r2.font.color.rgb = col if (len(it) == 3) else (NAVY if lvl == 0 else DARK)
        r2.font.name = FONT
    return tb


def add_image_fit(s, path, x, y, max_w, max_h, frame=True):
    """Place image scaled to fit max box, centered, optional border card."""
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = Emu(int(max_w / ar))
    else:
        h = max_h; w = Emu(int(max_h * ar))
    px = x + Emu(int((max_w - w) / 2))
    py = y + Emu(int((max_h - h) / 2))
    if frame:
        pad = Inches(0.06)
        card = rect(s, px - pad, py - pad, w + pad*2, h + pad*2, WHITE)
        card.line.color.rgb = LGREY; card.line.width = Pt(1.2)
    s.shapes.add_picture(path, px, py, width=w, height=h)


def kpi_card(s, x, y, w, h, value, label, vcolor=BLUE):
    c = rect(s, x, y, w, h, LGREY)
    c.line.color.rgb = WHITE
    rect(s, x, y, w, Pt(4), vcolor)
    txt(s, x, y + Inches(0.18), w, Inches(0.6),
        [(value, 30, True, vcolor)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.07), y + Inches(0.82), w - Inches(0.14), h - Inches(0.9),
        [(label, 11, False, GREY)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def table(s, x, y, w, rows, col_w, header_fill=BLUE, font=11, rh=0.34,
          header_font=11, zebra=True):
    nrows = len(rows); ncols = len(rows[0])
    gt = s.shapes.add_table(nrows, ncols, x, y, w, Inches(rh*nrows)).table
    # set column widths
    total = sum(col_w)
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Emu(int(w * cw / total))
    for i, row in enumerate(rows):
        gt.rows[i].height = Inches(rh)
        for j, cell in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.04)
            c.margin_top = Inches(0.01); c.margin_bottom = Inches(0.01)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            # cell may be (text, color, bold)
            if isinstance(cell, tuple):
                t, ccol, cb = cell
            else:
                t, ccol, cb = cell, None, False
            tfc = c.text_frame; tfc.word_wrap = True
            p = tfc.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(t); r.font.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(header_font)
            else:
                c.fill.solid()
                c.fill.fore_color.rgb = WHITE if (not zebra or i % 2 == 1) else LGREY
                r.font.color.rgb = ccol if ccol else DARK
                r.font.bold = cb; r.font.size = Pt(font)
    return gt


def divider(s, section_no, title, subtitle, idx):
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, Inches(3.05), SW, Pt(2.5), AMBER)
    rect(s, Inches(0.9), Inches(2.0), Inches(1.5), Inches(0.9), BLUE)
    txt(s, Inches(0.9), Inches(2.02), Inches(1.5), Inches(0.9),
        [(section_no, 40, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.7), Inches(2.0), Inches(9.8), Inches(1.0),
        [(title, 34, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.95), Inches(3.35), Inches(11.0), Inches(1.2),
        [(subtitle, 16, False, RGBColor(0xC7, 0xD3, 0xE3))])
    txt(s, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
        [(str(idx), 10, False, RGBColor(0x8A, 0x97, 0xA8))], align=PP_ALIGN.RIGHT)


_n = [0]
def nxt():
    _n[0] += 1
    return _n[0]


# ================================================================ SLIDE 1 — TITLE
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.28), SH, BLUE)
rect(s, Inches(0.28), 0, Inches(0.08), SH, TEAL)
rect(s, Inches(0.9), Inches(1.5), Inches(2.4), Pt(4), AMBER)
txt(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
    [("BENCHMARK SISTEMATICO · DEMAND FORECASTING NEL FRESH RETAIL", 14, True, TEAL)])
txt(s, Inches(0.85), Inches(2.35), Inches(11.7), Inches(1.9),
    [[("L'imputazione dello stockout conta\n", 40, True, WHITE)],
     [("davvero per il forecasting?", 40, True, WHITE)]], line_spacing=1.05)
txt(s, Inches(0.9), Inches(4.35), Inches(11.4), Inches(0.9),
    [("14 imputer × 8 forecaster = 113 combinazioni valutate su FreshRetailNet-50K "
      "(50.000 serie temporali, 898 negozi, 18 città).", 16, False, RGBColor(0xC7,0xD3,0xE3))],
    line_spacing=1.15)
# bottom bar
rect(s, 0, Inches(6.55), SW, Inches(0.95), RGBColor(0x0A, 0x17, 0x2B))
txt(s, Inches(0.9), Inches(6.7), Inches(8), Inches(0.6),
    [[("Francesco De Nicolo", 14, True, WHITE)],
     [("Framework statistico: Friedman χ² + Kendall's W + Nemenyi CD (Demšar 2006)", 11, False, RGBColor(0x9F,0xAE,0xC2))]])
txt(s, Inches(9.4), Inches(6.78), Inches(3.4), Inches(0.5),
    [("Giugno 2026", 13, True, AMBER)], align=PP_ALIGN.RIGHT)

# ================================================================ SLIDE 2 — EXECUTIVE SUMMARY
s = slide()
header(s, "Executive summary", "Il messaggio in una slide", nxt())
# big statement
box = rect(s, Inches(0.55), Inches(1.45), Inches(12.25), Inches(1.55), NAVY)
txt(s, Inches(0.85), Inches(1.55), Inches(11.7), Inches(1.35),
    [[("La famiglia di forecaster ", 17, False, WHITE), ("MLP_M5 domina il benchmark in ogni regime di volume", 17, True, AMBER),
      (". All'interno di MLP_M5 la scelta dell'imputer è ", 17, False, WHITE),
      ("praticamente irrilevante", 17, True, AMBER),
      (" — i lag features M5 disaccoppiano il forecasting dalla qualità della recovery. "
       "L'imputer conta solo per i forecaster senza lag.", 17, False, WHITE)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
# KPI row
y = Inches(3.35)
kpi_card(s, Inches(0.55), y, Inches(2.9), Inches(1.55), "itransformer\n__MLP_M5", "Best globale assoluto\n(mean rank 22.3 / 113)", BLUE)
kpi_card(s, Inches(3.65), y, Inches(2.9), Inches(1.55), "0.454", "Kendall's W globale\n(concordanza moderate)", TEAL)
kpi_card(s, Inches(6.75), y, Inches(2.9), Inches(1.55), "W ≈ 0", "Effetto imputer per\nMLP_M5 / LGB_M5", AMBER)
kpi_card(s, Inches(9.85), y, Inches(2.95), Inches(1.55), "26 / 113", "Celle Pareto-ottimali\n(trade-off accuracy/bias)", BLUE)
# takeaway strip
rect(s, Inches(0.55), Inches(5.25), Pt(5), Inches(1.55), AMBER)
bullets(s, Inches(0.8), Inches(5.25), Inches(12.0), Inches(1.7), [
    ("Per il practitioner: con lag features M5 disponibili, investire nel forecaster (MLP_M5), non nell'imputer.", 1),
    ("L'imputer ha valore solo per foundation models (Chronos/TimesFM) e naive aggregati, dove la recovery predice direttamente il forecasting.", 1),
    ("Trade-off strutturale accuracy↔bias: i modelli ML/DL sotto-stimano sistematicamente; per ridurre il bias servono naive aggregati.", 1),
], size=13.5, gap=6)

# ================================================================ SLIDE 3 — AGENDA
s = slide()
header(s, "Agenda", "Cosa copre questa presentazione", nxt())
items = [
    ("01", "Contesto & dataset", "Il problema del censoring da stockout · FreshRetailNet-50K", TEAL),
    ("02", "Disegno sperimentale", "Matrice 14×8 = 113 celle · metriche · framework statistico", BLUE),
    ("1", "RQ1 — Best cell & regimi di volume", "Best assoluto, Pareto, crossover, l'imputer aiuta?", AMBER),
    ("2", "RQ2 — La recovery predice il forecasting?", "Correlazione recovery → forecasting per famiglia", AMBER),
    ("3", "RQ3 — Foundation models per il retail", "Chronos-bolt e TimesFM: utili o no?", AMBER),
    ("★", "Sintesi & decision tree", "Messaggio scientifico + guida operativa per practitioner", BLUE),
    ("A", "Appendice tecnica", "Framework, matrice completa, tabelle, bibliografia", GREY),
]
y = Inches(1.45)
for no, t, sub, col in items:
    rect(s, Inches(0.7), y, Inches(0.72), Inches(0.62), col)
    txt(s, Inches(0.7), y, Inches(0.72), Inches(0.62), [(no, 18, True, WHITE)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.65), y - Inches(0.02), Inches(11), Inches(0.4), [(t, 16, True, NAVY)])
    txt(s, Inches(1.65), y + Inches(0.32), Inches(11), Inches(0.3), [(sub, 11.5, False, GREY)])
    y += Inches(0.76)

# ================================================================ SLIDE 4 — CONTESTO / PROBLEMA
s = slide()
header(s, "Contesto", "Il problema: vendite censurate dallo stockout", nxt())
rect(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(1.4), LGREY)
txt(s, Inches(0.8), Inches(1.55), Inches(5.6), Inches(1.2),
    [[("S_obs(t) = min( D(t), I(t) )", 20, True, NAVY)],
     [("Quando c'è stockout, le vendite osservate ≈ 0 a prescindere dalla domanda vera D(t).", 12.5, False, DARK)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
bullets(s, Inches(0.7), Inches(3.1), Inches(6.0), Inches(3.6), [
    ("Il censoring genera un ciclo vizioso", 0),
    ("I modelli imparano 'stockout → vendite 0' → sotto-stima sistematica della domanda latente.", 1),
    ("La soluzione classica: imputare i buchi", 0),
    ("Ricostruire la domanda nelle ore di stockout prima di addestrare il forecaster (pipeline two-stage).", 1),
    ("La domanda aperta", 0),
    ("Quale imputer scegliere? E soprattutto: quanto incide davvero la scelta dell'imputer sul forecasting finale?", 1),
], size=13.5, gap=6)
# right callout
rect(s, Inches(7.0), Inches(3.1), Inches(5.8), Inches(3.5), NAVY)
rect(s, Inches(7.0), Inches(3.1), Inches(5.8), Pt(4), AMBER)
txt(s, Inches(7.25), Inches(3.3), Inches(5.35), Inches(3.2),
    [[("Perché è rilevante per il business", 15, True, AMBER)],
     [("Retail deperibile (frutta, verdura, carne fresca): la sotto-stima della domanda porta a sotto-ordinare → stockout ricorrenti → vendite perse.", 13, False, WHITE)],
     [("", 6, False, WHITE)],
     [("Stimare correttamente la ", 13, False, WHITE), ("domanda latente", 13, True, AMBER),
      (" è la base per dimensionare il riassortimento.", 13, False, WHITE)]],
    line_spacing=1.12, space_after=8)

# ================================================================ SLIDE 5 — DATASET
s = slide()
header(s, "Dataset", "FreshRetailNet-50K — scala industriale, retail reale", nxt())
y = Inches(1.5)
cards = [("50.000", "serie temporali\n(store × SKU)"), ("898", "negozi in\n18 città cinesi"),
         ("865", "SKU deperibili\n(fresco)"), ("21,2%", "tasso medio\ndi stockout (ore 6–22)")]
x = Inches(0.55)
for v, l in cards:
    kpi_card(s, x, y, Inches(2.95), Inches(1.5), v, l, BLUE)
    x += Inches(3.07)
bullets(s, Inches(0.7), Inches(3.25), Inches(6.1), Inches(3.5), [
    ("Periodo: 28/03 → 02/07 2024 (97 giorni), granularità oraria.", 1),
    ("Orario operativo ristretto alle ore 6–22 (17 h/giorno): le ore notturne hanno vendite ~0 e distorcerebbero il WAPE.", 1),
    ("Split temporale: train gg 1–83 · validation gg 84–90 · test gg 91–97 (mai usato in training).", 1),
    ("Covariate shift nel test: temperatura +7 °C, pioggia +59% (passaggio primavera → estate).", 1),
    ("Distribuzione sbilanciata: City 0 = 52% del dataset.", 1),
], size=13, gap=8)
add_image_fit(s, f"{FIG}/fig_heatmap_general.png", Inches(7.0), Inches(3.2), Inches(5.9), Inches(3.6))
txt(s, Inches(7.0), Inches(6.85), Inches(5.9), Inches(0.3),
    [("Heatmap WAPE sulla matrice imputer × forecaster", 9, False, GREY)], align=PP_ALIGN.CENTER)

# ================================================================ SLIDE 6 — SETUP
s = slide()
header(s, "Disegno sperimentale", "La matrice: 14 imputer × 8 forecaster = 113 celle", nxt())
bullets(s, Inches(0.7), Inches(1.45), Inches(6.15), Inches(2.7), [
    ("14 imputer in 4 famiglie", 0),
    ("Baseline (no_imp) · naive aggregati (media/mediana × glob/cond) · classici (forward/seasonal/linear) · ML & DL (LGB, DLinear, SAITS, iTransformer, TimesNet, ImputeFormer, CSDI).", 1),
    ("8 forecaster in 4 famiglie", 0),
    ("Naive (Global/DoW Mean, MA K=21) · ML (LGB) · DL (MLP, TFT) · Foundation (Chronos-bolt, TimesFM).", 1),
], size=12.5, gap=6)
rect(s, Inches(7.05), Inches(1.45), Inches(5.75), Inches(2.55), LGREY)
rect(s, Inches(7.05), Inches(1.45), Inches(5.75), Pt(4), TEAL)
txt(s, Inches(7.3), Inches(1.6), Inches(5.3), Inches(2.35),
    [[("Metriche (solo ore in-stock del test)", 14, True, NAVY)],
     [("WAPE", 14, True, BLUE), (" = Σ|pred − obs| / Σobs   →  accuratezza", 13, False, DARK)],
     [("WPE", 14, True, BLUE), ("  = Σ(pred − obs) / Σobs   →  bias / sotto-stima", 13, False, DARK)],
     [("Con N=50K serie, p-value ≈ 0 ovunque → l'", 12, False, DARK),
      ("effect size", 12, True, TEAL), (" è la metrica discriminante.", 12, False, DARK)]],
    line_spacing=1.15, space_after=7)
# direct forecast + 3 RQ band
rect(s, Inches(0.55), Inches(4.25), Inches(12.25), Inches(0.95), NAVY)
txt(s, Inches(0.8), Inches(4.33), Inches(11.8), Inches(0.8),
    [[("Direct forecast", 13, True, AMBER),
      (": i lag features sono fissati all'anchor giorno 90 e applicati a tutti i 7 giorni di test. "
       "Variano solo le covariate esogene → questo isola l'effetto dell'imputation (l'unica differenza tra 'no_imp' e 'imputer X' è la qualità dei lag).",
       12.5, False, WHITE)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
# three RQ chips
rqs = [("RQ1", "Qual è la migliore combinazione e come cambia col volume? L'imputer aiuta?", AMBER),
       ("RQ2", "La qualità della recovery predice la qualità del forecasting?", TEAL),
       ("RQ3", "I foundation models sono competitivi sul retail deperibile?", BLUE)]
x = Inches(0.55)
for tag, d, col in rqs:
    rect(s, x, Inches(5.45), Inches(3.97), Inches(1.55), LGREY)
    rect(s, x, Inches(5.45), Inches(3.97), Pt(4), col)
    txt(s, x + Inches(0.15), Inches(5.55), Inches(3.7), Inches(0.4), [(tag, 16, True, col)])
    txt(s, x + Inches(0.15), Inches(5.95), Inches(3.7), Inches(1.0), [(d, 11.5, False, DARK)], line_spacing=1.08)
    x += Inches(4.1)

# ================================================================ SLIDE 7 — FRAMEWORK STATISTICO
s = slide()
header(s, "Metodo", "Un solo framework statistico, tre livelli di analisi", nxt())
steps = [
    ("Friedman χ²", "Rifiuta H0: 'tutte le k celle hanno la stessa distribuzione di rank'. Test non-parametrico su ranking cross-serie.", BLUE),
    ("Kendall's W", "Effect size globale ∈ [0,1]: quanto è concorde il ranking tra le serie. <0.1 negligible · 0.3–0.5 moderate · ≥0.5 large.", TEAL),
    ("Nemenyi CD", "Critical Difference post-hoc: due celle sono indistinguibili se |Δ mean-rank| ≤ CD → definisce l'equivalence set.", AMBER),
]
y = Inches(1.55)
for i, (t, d, col) in enumerate(steps):
    rect(s, Inches(0.7), y, Inches(12.1), Inches(1.15), LGREY)
    rect(s, Inches(0.7), y, Pt(6), Inches(1.15), col)
    txt(s, Inches(1.0), y + Inches(0.12), Inches(3.0), Inches(0.9), [(t, 18, True, col)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.0), y + Inches(0.1), Inches(8.6), Inches(0.95), [(d, 13, False, DARK)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    if i < 2:
        txt(s, Inches(6.4), y + Inches(1.02), Inches(0.6), Inches(0.3), [("▼", 12, True, GREY)], align=PP_ALIGN.CENTER)
    y += Inches(1.42)
rect(s, Inches(0.7), Inches(5.95), Inches(12.1), Inches(0.95), NAVY)
txt(s, Inches(0.95), Inches(6.03), Inches(11.6), Inches(0.8),
    [[("Stesso framework applicato a 3 scope:  ", 13, True, AMBER),
      ("(1) globale sulla matrice intera · (2) stratificato per quartile di volume · (3) ristretto per singolo forecaster. "
       "Cliff's δ resta come effect size descrittivo (mai come decision rule).", 12.5, False, WHITE)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)

# ================================================================ DIVIDER RQ1
s = slide()
divider(s, "1", "Best cell, regimi di volume\ne ruolo dell'imputer",
        "RQ1 — Qual è la combinazione vincente? Come cambia col volume? L'imputazione aiuta davvero?", nxt())

# ---------------------------------------------------------------- 1.1 best globale
s = slide()
header(s, "RQ1.1 — Best globale", "Una sola famiglia di forecaster vince: MLP_M5", nxt())
bullets(s, Inches(0.7), Inches(1.4), Inches(6.0), Inches(2.6), [
    ("Best assoluto: itransformer__MLP_M5", 0),
    ("Mean rank 22.27 su 113 celle. Kendall's W = 0.454 (moderate) → il ranking è generalizzabile cross-serie.", 1),
    ("Equivalence set di sole 2 celle", 0),
    ("itransformer__MLP_M5 (rank 22.27) e lgb__MLP_M5 (rank 22.78, Δ=0.51 ≤ CD=0.903).", 1),
    ("Entrambe sono MLP_M5", 0, AMBER),
    ("La famiglia di forecaster vincente è isolata: l'imputer cambia, la famiglia no.", 1),
], size=13, gap=6)
add_image_fit(s, f"{FIG}/fig_cd_diagram.png", Inches(6.9), Inches(1.45), Inches(6.0), Inches(4.0))
txt(s, Inches(6.9), Inches(5.55), Inches(6.0), Inches(0.3),
    [("Critical Difference diagram — ranking globale delle 113 celle", 9, False, GREY)], align=PP_ALIGN.CENTER)
rect(s, Inches(0.7), Inches(5.95), Inches(6.0), Inches(0.95), LGREY)
rect(s, Inches(0.7), Inches(5.95), Pt(5), Inches(0.95), AMBER)
txt(s, Inches(0.95), Inches(6.03), Inches(5.6), Inches(0.8),
    [[("Finding 1.1: ", 13, True, NAVY), ("il vincitore è la famiglia MLP_M5, non un imputer specifico. W moderate → risultato robusto.", 12.5, False, DARK)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

# ---------------------------------------------------------------- 1.1b Pareto globale
s = slide()
header(s, "RQ1.1 — Trade-off accuracy × bias", "Pareto frontier globale: accuracy e bias non si massimizzano insieme", nxt())
add_image_fit(s, f"{FIG}/fig19_pareto_frontier.png", Inches(0.6), Inches(1.45), Inches(6.6), Inches(5.0))
txt(s, Inches(0.6), Inches(6.55), Inches(6.6), Inches(0.3),
    [("Pareto su (WAPE, |WPE|) — 26/113 celle non dominate", 9, False, GREY)], align=PP_ALIGN.CENTER)
table(s, Inches(7.4), Inches(1.6), Inches(5.5), [
    ["Ruolo", "Cella", "WAPE", "|WPE|"],
    ["Best WAPE", "timesnet__MLP_M5", ("0.973", BLUE, True), "0.886"],
    ["Knee (bilanciato)", "mediana_glob__dow_mean", "1.101", ("0.190", TEAL, True)],
    ["Min |WPE|", "linear_interp__timesfm", "1.295", ("0.061", TEAL, True)],
], [1.5, 2.2, 0.9, 0.9], font=10.5, header_font=10.5, rh=0.5)
bullets(s, Inches(7.4), Inches(3.5), Inches(5.5), Inches(3.3), [
    ("Il trade-off è strutturale.", 0, AMBER),
    ("Il best-WAPE ha sempre |WPE| ≥ 0.77: i forecaster ML/DL sotto-stimano sistematicamente.", 1),
    ("Per ridurre il bias servono naive aggregati (Global/DoW/MA), che pagano in WAPE.", 1),
    ("Il knee point mediana_glob__dow_mean è il compromesso ragionevole per chi valuta sia accuracy che bias.", 1),
], size=12.5, gap=6)

# ---------------------------------------------------------------- 1.2a best per volume
s = slide()
header(s, "RQ1.2 — Robustezza per volume", "Crossover soft: MLP_M5 vince ovunque, cambia solo l'imputer", nxt())
table(s, Inches(0.6), Inches(1.5), Inches(12.2), [
    ["Quartile", "Range volume", "Best cell", "Kendall W", "CD-equiv"],
    ["Q1 (basso)", "[11, 40]", ("lgb__MLP_M5", BLUE, True), "0.653 (large)", "12"],
    ["Q2", "(40, 54]", ("lgb__MLP_M5", BLUE, True), "0.586 (large)", "13"],
    ["Q3 (medio-alto)", "(54, 86]", ("itransformer__MLP_M5", BLUE, True), "0.417 (moderate)", "4"],
    ["Q4 (alto)", "(86, 5326]", ("itransformer__MLP_M5", BLUE, True), "0.396 (moderate)", "2"],
], [1.5, 1.6, 2.6, 1.8, 1.2], font=12.5, header_font=12, rh=0.46)
bullets(s, Inches(0.7), Inches(4.0), Inches(7.2), Inches(2.9), [
    ("La famiglia MLP_M5 è invariata in tutti i regimi.", 0, AMBER),
    ("Cambia solo l'imputer ottimale: lgb in basso volume (Q1/Q2), itransformer in alto volume (Q3/Q4).", 1),
    ("Q4 è il regime più discriminante (solo 2 celle equivalenti); Q1 il più saturato (12).", 1),
    ("W e numero di equivalenti non sono ridondanti: misurano concordanza globale e separazione locale dentro la famiglia vincente.", 1),
], size=12.5, gap=6)
rect(s, Inches(8.2), Inches(4.0), Inches(4.6), Inches(2.7), NAVY)
rect(s, Inches(8.2), Inches(4.0), Inches(4.6), Pt(4), AMBER)
txt(s, Inches(8.45), Inches(4.2), Inches(4.15), Inches(2.4),
    [[("In pratica", 14, True, AMBER)],
     [("La scelta del forecaster (MLP_M5) è stabile a prescindere dal volume della serie.", 13, False, WHITE)],
     [("La scelta dell'imputer è un dettaglio di second'ordine che dipende solo dal volume.", 13, False, WHITE)]],
    line_spacing=1.13, space_after=9)

# ---------------------------------------------------------------- 1.2b pareto per quartile
s = slide()
header(s, "RQ1.2 — Pareto per quartile", "Più volume = meno celle ottimali = scelta più decisiva", nxt())
add_image_fit(s, f"{FIG}/fig20_pareto_per_quartile.png", Inches(0.55), Inches(1.45), Inches(7.2), Inches(5.05))
txt(s, Inches(0.55), Inches(6.55), Inches(7.2), Inches(0.3),
    [("Pareto frontier WAPE × |WPE| dentro ciascun quartile", 9, False, GREY)], align=PP_ALIGN.CENTER)
table(s, Inches(8.0), Inches(1.55), Inches(4.9), [
    ["Q", "# Pareto", "Top-WAPE"],
    ["Q1", ("28", AMBER, True), "1.000"],
    ["Q2", "22", "0.995"],
    ["Q3", "15", "0.958"],
    ["Q4", ("12", AMBER, True), ("0.757", BLUE, True)],
], [0.7, 1.1, 1.2], font=12, header_font=11, rh=0.42)
bullets(s, Inches(8.0), Inches(3.5), Inches(4.9), Inches(3.3), [
    ("# Pareto cala 28→12 da Q1 a Q4: alto volume = più discriminante.", 1),
    ("Basso volume: TFT popola la coda low-bias.", 1),
    ("Alto volume: MLP_M5 al top-WAPE, naive aggregati al low-bias — in Q4 i naive battono i ML su |WPE| (0.06 vs 0.41) pagando solo +3% WAPE.", 1),
    ("Chronos appare solo in Q2; TimesFM mai nei quartili.", 1),
], size=12, gap=6)

# ---------------------------------------------------------------- 1.2c crossover
s = slide()
header(s, "RQ1.2 — Crossover per volume", "L'alto volume comprime le differenze tra famiglie", nxt())
add_image_fit(s, f"{FIG}/fig_rq3_crossover_fixed_global.png", Inches(0.55), Inches(1.45), Inches(7.0), Inches(4.4))
txt(s, Inches(0.55), Inches(5.9), Inches(7.0), Inches(0.3),
    [("Evoluzione WAPE per famiglia di forecaster, Q1 → Q4", 9, False, GREY)], align=PP_ALIGN.CENTER)
bullets(s, Inches(7.75), Inches(1.5), Inches(5.1), Inches(3.0), [
    ("Convergenza in Q4: ML, TFT e naive tutti a WAPE ≈ 0.77.", 1),
    ("Naive vs ML: il gap scende da +25% (Q1) a 0% (Q4).", 1),
    ("Chronos resta piatto (~1.00): non sfrutta il volume → perde relativamente.", 1),
    ("TimesFM è il peggiore in tutti i quartili.", 1),
], size=12.5, gap=7)
rect(s, Inches(7.75), Inches(4.5), Inches(5.1), Inches(2.35), NAVY)
rect(s, Inches(7.75), Inches(4.5), Inches(5.1), Pt(4), AMBER)
txt(s, Inches(8.0), Inches(4.62), Inches(4.6), Inches(2.15),
    [[("Decision tree practitioner", 13, True, AMBER)],
     [("Basso volume:", 12.5, True, WHITE), (" ML con lag M5 (MLP/LGB/TFT); naive nettamente peggiori.", 12.5, False, WHITE)],
     [("Alto volume:", 12.5, True, WHITE), (" la famiglia conta poco; i naive sono competitivi con bias molto migliore.", 12.5, False, WHITE)],
     [("Foundation:", 12.5, True, WHITE), (" solo baseline zero-shot, mai dominanti.", 12.5, False, WHITE)]],
    line_spacing=1.08, space_after=5)

# ---------------------------------------------------------------- 1.3 imputer aiuta?
s = slide()
header(s, "RQ1.3 — L'imputer aiuta?", "Dicotomia netta: dipende dall'architettura del forecaster", nxt())
add_image_fit(s, f"{FIG}/fig_cd_per_forecaster.png", Inches(0.55), Inches(1.45), Inches(6.5), Inches(5.0))
txt(s, Inches(0.55), Inches(6.55), Inches(6.5), Inches(0.3),
    [("CD ristretto: posizione di 'no_imp' per ciascun forecaster", 9, False, GREY)], align=PP_ALIGN.CENTER)
table(s, Inches(7.25), Inches(1.5), Inches(5.6), [
    ["Forecaster", "W", "Imputer aiuta?"],
    ["Global / DoW / MA", "0.44–0.47", ("SÌ", GREEN, True)],
    ["TimesFM", "0.174", ("SÌ", GREEN, True)],
    ["Chronos-bolt", "0.222", ("SÌ", GREEN, True)],
    ["TFT", "0.220", ("SÌ", GREEN, True)],
    ["MLP_M5", "0.029", ("NO", RED, True)],
    ["LGB_M5", "0.009", ("NO", RED, True)],
], [2.2, 1.3, 1.6], font=11.5, header_font=11, rh=0.42)
bullets(s, Inches(7.25), Inches(4.7), Inches(5.6), Inches(2.2), [
    ("Imputer-sensitive (naive, foundation, TFT): l'imputer aiuta sempre, ranking generalizzabile (W ≥ 0.22).", 1),
    ("Imputer-irrelevant (MLP_M5, LGB_M5): W ≈ 0 → ranking imputer praticamente casuale; 'no_imp' è equivalente al best.", 1),
], size=12, gap=6)

# ---------------------------------------------------------------- 1.3.1 loss sensitivity
s = slide()
header(s, "RQ1.3 — Robustezza alla loss", "Il finding 'l'imputer non conta' sopravvive al cambio di loss", nxt())
bullets(s, Inches(0.7), Inches(1.45), Inches(5.6), Inches(3.2), [
    ("Il dubbio", 0),
    ("'W ≈ 0 per MLP_M5/LGB_M5' è un artefatto della loss MSE non allineata al WAPE?", 1),
    ("Il test", 0),
    ("Ri-addestrate 24 celle con loss MAE; confronto pairwise Wilcoxon + Cliff δ.", 1),
    ("Il risultato", 0, AMBER),
    ("MAE migliora il WAPE in tutte le 24 celle, ma lo spread inter-imputer resta piccolo sotto entrambe le loss.", 1),
    ("Conclusione: il finding è strutturale, non causato dalla loss.", 1, GREEN),
], size=12.5, gap=5)
add_image_fit(s, f"{FIG}/fig_rq5_mae_vs_mse.png", Inches(6.6), Inches(1.5), Inches(6.3), Inches(4.9))
txt(s, Inches(6.6), Inches(6.5), Inches(6.3), Inches(0.3),
    [("Sensitivity analysis: training loss MAE vs MSE", 9, False, GREY)], align=PP_ALIGN.CENTER)

# ================================================================ DIVIDER RQ2
s = slide()
divider(s, "2", "La recovery predice\nil forecasting?",
        "RQ2 — Una buona ricostruzione dei dati mancanti garantisce un buon forecasting downstream?", nxt())

# ---------------------------------------------------------------- 2
s = slide()
header(s, "RQ2 — Recovery → Forecasting", "Sì, ma solo quanto il forecaster dipende dai dati imputati", nxt())
add_image_fit(s, f"{FIG}/fig_rq2_recovery_vs_forecasting.png", Inches(0.55), Inches(1.45), Inches(6.4), Inches(5.0))
txt(s, Inches(0.55), Inches(6.55), Inches(6.4), Inches(0.3),
    [("Correlazione per-serie recovery ↔ forecasting (Cliff δ vs 0)", 9, False, GREY)], align=PP_ALIGN.CENTER)
table(s, Inches(7.2), Inches(1.5), Inches(5.65), [
    ["Famiglia", "Cliff δ", "Categoria"],
    ["Naive (MA/DoW/Global)", ("+0.81…+0.87", GREEN, True), "LARGE"],
    ["Foundation (TimesFM/Chronos)", ("+0.47…+0.56", TEAL, True), "medium–LARGE"],
    ["ML+lag (MLP/LGB M5)", ("+0.14…+0.18", GREY, True), "small"],
    ["TFT (DL+lag)", ("−0.24", RED, True), "small (inverso)"],
], [2.6, 1.5, 1.6], font=11, header_font=11, rh=0.46)
bullets(s, Inches(7.2), Inches(3.7), Inches(5.65), Inches(3.1), [
    ("Naive: la predizione è quasi una funzione lineare dei valori imputati → recovery → forecasting deterministico.", 1),
    ("Foundation: ricevono il segnale ma lo trasformano → effetto attenuato ma reale.", 1),
    ("ML/DL con lag M5: i lag fanno da buffer → decoupling dalla qualità dell'imputer.", 1),
    ("Chiave del Finding 1.3: i lag M5 isolano il forecasting dalla recovery.", 1, AMBER),
], size=11.5, gap=6)

# ================================================================ DIVIDER RQ3
s = slide()
divider(s, "3", "Foundation models\nper il retail",
        "RQ3 — Chronos-bolt e TimesFM sono competitivi sul retail deperibile?", nxt())

# ---------------------------------------------------------------- 3
s = slide()
header(s, "RQ3 — Foundation models", "Recovery-sensitive ma dominati da MLP_M5", nxt())
y = Inches(1.5)
kpi_card(s, Inches(0.55), y, Inches(3.9), Inches(1.5), "1.007", "Chronos-bolt × no_imp\nWAPE (competitivo)", TEAL)
kpi_card(s, Inches(4.65), y, Inches(3.9), Inches(1.5), "+18%", "TimesFM peggio di Chronos\n(CPU only, 5× più lento)", AMBER)
kpi_card(s, Inches(8.75), y, Inches(4.05), Inches(1.5), "~35–40", "Mean rank foundation\nvs 22.3 del best MLP_M5", RED)
bullets(s, Inches(0.7), Inches(3.3), Inches(12.0), Inches(2.4), [
    ("Entrambi recovery-sensitive", 0),
    ("Cliff δ vs 0: TimesFM +0.556 (LARGE), Chronos +0.472 (medium). Best imputer coerente per entrambi: imputeformer.", 1),
    ("Ma dominati da MLP_M5 sulla matrice principale", 0, AMBER),
    ("Mean rank ≈ 35–40 contro 22.3 di itransformer__MLP_M5. Nessun regime di volume in cui dominano.", 1),
    ("Ruolo pratico", 0),
    ("Utili come baseline zero-shot quando non si hanno lag features; non competitivi sul retail deperibile con lag disponibili.", 1),
], size=13, gap=6)
rect(s, Inches(0.7), Inches(6.15), Inches(12.1), Inches(0.75), NAVY)
txt(s, Inches(0.95), Inches(6.22), Inches(11.6), Inches(0.6),
    [[("Finding 3: ", 13, True, AMBER), ("i foundation models pagano la mancanza di lag features; il segnale imputato li influenza ma non basta a colmare il gap con MLP_M5.", 12.5, False, WHITE)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

# ================================================================ DIVIDER SINTESI
s = slide()
divider(s, "★", "Sintesi & guida operativa",
        "Il messaggio scientifico e cosa farne in pratica", nxt())

# ---------------------------------------------------------------- messaggio chiave
s = slide()
header(s, "Sintesi", "Il messaggio scientifico chiave", nxt())
rect(s, Inches(0.7), Inches(1.5), Inches(12.1), Inches(2.4), NAVY)
rect(s, Inches(0.7), Inches(1.5), Pt(6), Inches(2.4), AMBER)
txt(s, Inches(1.05), Inches(1.7), Inches(11.5), Inches(2.05),
    [[("La famiglia MLP_M5 domina il benchmark in ogni regime. ", 18, True, WHITE),
      ("Dentro MLP_M5 la scelta dell'imputer è praticamente irrilevante (W ≈ 0) perché i lag features M5 disaccoppiano il forecasting dalla recovery quality. "
       "L'imputer conta solo per i forecaster senza lag — foundation models e naive aggregati — dove la recovery predice direttamente il forecasting (Cliff δ ≥ +0.47).",
       16, False, RGBColor(0xDD,0xE5,0xF0))]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
cols = [("Best globale", "itransformer__MLP_M5 · 2 celle equivalenti · W=0.454 moderate", BLUE),
        ("Robustezza", "MLP_M5 vince in Q1→Q4; cambia solo l'imputer (lgb→itransformer)", TEAL),
        ("Imputer", "Irrilevante con lag M5; decisivo per naive & foundation", AMBER)]
x = Inches(0.7)
for t, d, col in cols:
    rect(s, x, Inches(4.25), Inches(3.93), Inches(2.45), LGREY)
    rect(s, x, Inches(4.25), Inches(3.93), Pt(4), col)
    txt(s, x + Inches(0.2), Inches(4.4), Inches(3.55), Inches(0.5), [(t, 15, True, col)])
    txt(s, x + Inches(0.2), Inches(4.95), Inches(3.55), Inches(1.6), [(d, 13, False, DARK)], line_spacing=1.15)
    x += Inches(4.06)

# ---------------------------------------------------------------- decision tree
s = slide()
header(s, "Guida operativa", "Decision tree per il practitioner", nxt())
rows = [
    ("Hai lag features (storico ≥ 14 gg)?", "Usa MLP_M5 (o LGB_M5). L'imputer è secondario: 'no_imp' va benissimo.", GREEN),
    ("Basso volume (Q1–Q2)?", "ML con lag M5 (MLP/LGB/TFT). I naive sono nettamente peggiori. Foundation solo come baseline.", BLUE),
    ("Alto volume (Q4)?", "La famiglia conta poco: i naive aggregati sono competitivi e con bias molto migliore (|WPE| 0.06 vs 0.41).", TEAL),
    ("Ti serve basso bias (no sotto-stima)?", "Naive aggregati (Global/DoW/MA) o il knee point mediana_glob__dow_mean. Paghi un po' di WAPE.", AMBER),
    ("Niente lag / cold start?", "Foundation zero-shot (Chronos-bolt) + imputer imputeformer. Baseline ragionevole, non ottimale.", GREY),
]
y = Inches(1.5)
for q, a, col in rows:
    rect(s, Inches(0.7), y, Inches(4.45), Inches(0.95), col)
    txt(s, Inches(0.85), y, Inches(4.2), Inches(0.95), [(q, 13.5, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.5), y + Inches(0.06), Inches(0.5), Inches(0.85), [("→", 18, True, col)], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(6.0), y, Inches(6.8), Inches(0.95), LGREY)
    txt(s, Inches(6.2), y, Inches(6.45), Inches(0.95), [(a, 12.5, False, DARK)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
    y += Inches(1.07)

# ---------------------------------------------------------------- limitazioni
s = slide()
header(s, "Trasparenza", "Limitazioni dichiarate", nxt())
bullets(s, Inches(0.7), Inches(1.5), Inches(6.05), Inches(5.2), [
    ("Un solo dataset", 0),
    ("Retail cinese deperibile, 3.5 mesi: niente stagionalità annuale.", 1),
    ("Orizzonte breve", 0),
    ("7 giorni di test; orizzonti più lunghi non valutati.", 1),
    ("HP non ottimizzati per cella", 0),
    ("HP standardizzati (sensitivity HPO in appendice mostra effetto piccolo).", 1),
], size=13, gap=6)
bullets(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.2), [
    ("Range stockout ristretto", 0),
    ("16–31% per il 90% delle serie; stockout estremi (>50%) non validati.", 1),
    ("Sbilanciamento geografico", 0),
    ("City 0 = 52% del dataset (mitigato da modello globale + city_id).", 1),
    ("Asimmetrie di pipeline", 0),
    ("Alcuni imputer DL slim per vincoli CPU; 4 celle TFT escluse per OOM.", 1),
], size=13, gap=6)

# ---------------------------------------------------------------- conclusioni
s = slide()
header(s, "Conclusioni", "Cosa portare a casa", nxt())
points = [
    ("1", "MLP_M5 è la scelta di default", "Vince in ogni regime di volume con ranking generalizzabile (W moderate). È la decisione che conta di più.", BLUE),
    ("2", "L'imputer è sopravvalutato (con lag)", "Con lag features M5 + loss allineata, la scelta dell'imputer è rumore. Risparmia complessità di pipeline.", TEAL),
    ("3", "L'imputer conta senza lag", "Per naive e foundation models la recovery predice il forecasting: lì investire in un buon imputer (imputeformer/mediana_glob) ripaga.", AMBER),
    ("4", "Accuracy e bias sono in trade-off", "Nessuna singola soluzione ottimale: scegliere lungo la Pareto frontier in base al costo relativo di errore vs sotto-stima.", GREEN),
]
y = Inches(1.5)
for no, t, d, col in points:
    rect(s, Inches(0.7), y, Inches(0.85), Inches(1.18), col)
    txt(s, Inches(0.7), y, Inches(0.85), Inches(1.18), [(no, 26, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.7), y, Inches(11.1), Inches(1.18), LGREY)
    txt(s, Inches(1.95), y + Inches(0.12), Inches(10.6), Inches(0.4), [(t, 16, True, NAVY)])
    txt(s, Inches(1.95), y + Inches(0.55), Inches(10.6), Inches(0.6), [(d, 12.5, False, DARK)], line_spacing=1.05)
    y += Inches(1.32)

# ================================================================ DIVIDER APPENDICE
s = slide()
divider(s, "A", "Appendice tecnica",
        "Framework dettagliato, matrice completa, tabelle numeriche, bibliografia", nxt())

# ---------------------------------------------------------------- A1 framework dettaglio
s = slide()
header(s, "Appendice A1", "Framework statistico — mappa completa", nxt())
table(s, Inches(0.6), Inches(1.5), Inches(12.25), [
    ["Sezione", "Domanda", "Test", "Effect size"],
    ["RQ1.1 / 1.2 / 1.3", "Best cell + equivalence set (k>2 metodi)", "Friedman χ²", "Kendall's W + Nemenyi CD"],
    ["RQ1.1b / 1.2b", "Trade-off WAPE × |WPE|", "dominance", "n. celle Pareto-ottimali"],
    ["RQ1.2c", "Crossover famiglie per quartile", "descrittivo", "Δ WAPE per famiglia"],
    ["RQ1.3.1", "Robustezza loss MAE vs MSE", "Wilcoxon paired", "Cliff δ"],
    ["RQ2", "Recovery vs forecasting (correlazione)", "Wilcoxon vs 0 (su ρ_i)", "Cliff δ vs 0 + CI bootstrap"],
    ["RQ3", "Foundation models", "come RQ1 + RQ2", "come RQ1 + RQ2"],
], [1.8, 3.6, 2.2, 2.6], font=11.5, header_font=11.5, rh=0.5)
txt(s, Inches(0.6), Inches(5.7), Inches(12.2), Inches(1.0),
    [[("Nota metodologica: ", 12.5, True, NAVY),
      ("nessun TOST e nessuna soglia Cliff δ < 0.147 come decision rule. Cliff δ è usato solo come effect size descrittivo. "
       "Gli script TOST/threshold restano come materiale supplementare non citato nel paper.", 12, False, DARK)]],
    line_spacing=1.12)

# ---------------------------------------------------------------- A2 matrice imputer/forecaster
s = slide()
header(s, "Appendice A2", "La matrice completa: 14 imputer × 8 forecaster", nxt())
table(s, Inches(0.55), Inches(1.45), Inches(6.1), [
    ["Famiglia imputer", "Imputer"],
    ["Baseline", "no_imp"],
    ["Naive aggregati", "media/mediana × glob/cond (4)"],
    ["Classici TS", "forward_fill, seasonal_naive, linear_interp"],
    ["ML / DL", "lgb, dlinear, saits"],
    ["SOTA (paper)", "iTransformer, TimesNet, ImputeFormer, CSDI"],
], [2.0, 3.2], font=11.5, header_font=11.5, rh=0.55)
table(s, Inches(6.95), Inches(1.45), Inches(5.85), [
    ["Famiglia forecaster", "Forecaster"],
    ["Naive", "Global Mean, DoW Mean, MA (K=56)"],
    ["ML tabellare", "LGB (M5 lags)"],
    ["Deep Learning", "MLP (M5 lags), TFT"],
    ["Foundation", "Chronos-bolt, TimesFM 2.5"],
], [2.0, 3.0], font=11.5, header_font=11.5, rh=0.55)
rect(s, Inches(0.55), Inches(5.2), Inches(12.25), Inches(1.5), LGREY)
rect(s, Inches(0.55), Inches(5.2), Pt(5), Inches(1.5), TEAL)
txt(s, Inches(0.85), Inches(5.35), Inches(11.7), Inches(1.25),
    [[("Lag M5-style (11 feature): ", 13, True, NAVY),
      ("lag_1d, lag_7d, lag_14d, rmean_7d, rmean_14d, rstd_7d, lag_dow, rmean_dow, daily_total_lag1, daily_total_rmean7, momentum_1d_7d.", 12.5, False, DARK)],
     [("Sono questi lag a fare da 'buffer' che disaccoppia MLP_M5 / LGB_M5 dalla qualità dell'imputer.", 12.5, False, GREY, True)]],
    line_spacing=1.15, space_after=6)

# ---------------------------------------------------------------- A3 tabella 1.3 completa
s = slide()
header(s, "Appendice A3", "RQ1.3 — L'imputer aiuta? Tabella completa", nxt())
table(s, Inches(0.7), Inches(1.5), Inches(12.0), [
    ["Forecaster", "k", "Friedman best", "no_imp pos.", "Kendall W", "Categoria", "Aiuta?"],
    ["Global Mean", "14", "mediana_glob", "6°/14", "0.469", "moderate", ("SÌ", GREEN, True)],
    ["MA_K56", "14", "mediana_glob", "6°/14", "0.463", "moderate", ("SÌ", GREEN, True)],
    ["DoW Mean", "14", "mediana_glob", "5°/14", "0.445", "moderate", ("SÌ", GREEN, True)],
    ["TimesFM", "14", "imputeformer", "2°/14", "0.174", "small", ("SÌ", GREEN, True)],
    ["Chronos-bolt", "14", "imputeformer", "2°/14", "0.222", "small", ("SÌ", GREEN, True)],
    ["TFT", "13", "dlinear", "2°/13", "0.220", "small", ("SÌ", GREEN, True)],
    ["MLP_M5", "14", "itransformer", "9°/14", "0.029", "negligible", ("NO", RED, True)],
    ["LGB_M5", "14", "mediana_glob", "8°/14", "0.009", "negligible", ("NO", RED, True)],
], [1.9, 0.5, 1.9, 1.2, 1.2, 1.4, 1.0], font=11, header_font=10.5, rh=0.43)
txt(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.7),
    [("Best imputer coerente: imputeformer per foundation models · mediana_glob per naive · dlinear per TFT. "
      "Per MLP_M5 / LGB_M5 il ranking degli imputer è praticamente casuale (W ≈ 0).", 12, False, GREY)], line_spacing=1.1)

# ---------------------------------------------------------------- A4 recovery table
s = slide()
header(s, "Appendice A4", "RQ2 — Recovery → Forecasting, dettaglio per forecaster", nxt())
table(s, Inches(1.0), Inches(1.5), Inches(11.3), [
    ["Forecaster", "Famiglia", "median ρ", "Cliff δ vs 0", "Categoria Romano", "Predice?"],
    ["MA_K56", "naive", "+0.758", ("+0.847", GREEN, True), "LARGE", "SÌ"],
    ["DoW Mean", "naive", "+0.867", ("+0.823", GREEN, True), "LARGE", "SÌ"],
    ["Global Mean", "naive", "+0.883", ("+0.813", GREEN, True), "LARGE", "SÌ"],
    ["TimesFM", "foundation", "+0.333", ("+0.556", TEAL, True), "LARGE", "SÌ"],
    ["Chronos-bolt", "foundation", "+0.367", ("+0.472", TEAL, True), "medium", "SÌ"],
    ["MLP_M5", "ML+lag", "+0.100", ("+0.183", GREY, True), "small", "debole"],
    ["LGB_M5", "ML+lag", "+0.100", ("+0.142", GREY, True), "borderline", "quasi no"],
    ["TFT", "DL+lag", "−0.143", ("−0.239", RED, True), "small (inverso)", "no, opposto"],
], [1.7, 1.4, 1.2, 1.4, 1.9, 1.5], font=11, header_font=10.5, rh=0.42)
txt(s, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.6),
    [("La dipendenza recovery→forecasting cresce con la dipendenza diretta dai dati imputati: naive ≫ foundation ≫ ML+lag.", 12, False, GREY)],
    align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- A5 recovery quality (traccia A)
s = slide()
header(s, "Appendice A5", "Qualità dell'imputation (Traccia A — MNAR recovery)", nxt())
add_image_fit(s, f"{FIG}/fig_rq2_per_series_spearman.png", Inches(7.0), Inches(1.5), Inches(5.9), Inches(5.0))
txt(s, Inches(7.0), Inches(6.55), Inches(5.9), Inches(0.3),
    [("Distribuzione per-serie della correlazione di Spearman", 9, False, GREY)], align=PP_ALIGN.CENTER)
table(s, Inches(0.6), Inches(1.5), Inches(6.1), [
    ["Imputer", "WAPE_rec", "WPE_rec"],
    ["Mediana globale", ("0.809", GREEN, True), "−0.57"],
    ["Mediana condizionata", "0.846", "−0.47"],
    ["ImputeFormer", "0.867", "−0.74"],
    ["LGB imputer", "0.930", "−0.10"],
    ["iTransformer", "0.930", "−0.45"],
    ["SAITS", "0.943", "−0.82"],
    ["DLinear", "0.951", "−0.74"],
    ["TimesNet", "1.041", "−0.87"],
    ["Forward Fill", ("1.188", RED, True), "+0.13"],
], [2.3, 1.3, 1.2], font=10.5, header_font=10.5, rh=0.40)
txt(s, Inches(0.6), Inches(5.95), Inches(6.1), Inches(0.9),
    [("La Mediana globale ha la migliore recovery; ma — Finding chiave — la qualità della recovery NON è proxy del forecasting downstream (TimesNet, peggior recovery, dà il miglior MLP_M5).",
      11, False, GREY)], line_spacing=1.12)

# ---------------------------------------------------------------- A6 bibliografia
s = slide()
header(s, "Appendice A6", "Riferimenti bibliografici chiave", nxt())
refs = [
    ("Liu et al. (2025)", "FreshRetailNet-50K: Latent Demand from 50,000 Stores for World-scale Stockout Prediction in Fresh Retail. arXiv:2505.16319."),
    ("Demšar (2006)", "Statistical comparisons of classifiers over multiple data sets. JMLR. — Friedman + Nemenyi CD."),
    ("Romano et al. (2006)", "Appropriate statistics for ordinal level data. — soglie Cliff δ."),
    ("Du et al. (2023)", "SAITS: Self-Attention-based Imputation for Time Series. NeurIPS."),
    ("Du et al. (2023)", "PyPOTS: A Python Toolbox for Data Mining on Partially-Observed Time Series. arXiv:2305.18811."),
    ("Zeng et al. (2022)", "DLinear: Are Transformers Effective for Time Series Forecasting? AAAI 2023."),
    ("Ansari et al. (2024)", "Chronos: Learning the Language of Time Series. arXiv:2403.07815."),
]
y = Inches(1.55)
for a, t in refs:
    rect(s, Inches(0.7), y + Inches(0.07), Pt(4), Inches(0.55), BLUE)
    txt(s, Inches(0.95), y, Inches(3.0), Inches(0.7), [(a, 13, True, NAVY)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.1), y, Inches(8.7), Inches(0.7), [(t, 12, False, DARK)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    y += Inches(0.74)

# ---------------------------------------------------------------- CLOSING
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.0), SW, Pt(2.5), AMBER)
rect(s, Inches(0.9), Inches(2.1), Inches(2.4), Pt(4), TEAL)
txt(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.0),
    [("Grazie.", 40, True, WHITE)])
txt(s, Inches(0.95), Inches(3.35), Inches(11.4), Inches(1.4),
    [[("Domanda di ricerca, in una riga: ", 16, True, AMBER)],
     [("con lag features disponibili, investi nel forecaster (MLP_M5) — non nell'imputer. "
       "L'imputazione conta solo quando il forecaster non ha lag.", 16, False, RGBColor(0xC7,0xD3,0xE3))]],
    line_spacing=1.2)
txt(s, Inches(0.95), Inches(5.4), Inches(11.4), Inches(0.5),
    [("FreshRetailNet-50K · 14 imputer × 8 forecaster · 113 celle · Friedman + Kendall's W + Nemenyi CD", 12, False, RGBColor(0x9F,0xAE,0xC2))])

# ---------------------------------------------------------------- save
out = "FreshRetailNet_Imputer_Forecaster_Presentazione.pptx"
prs.save(out)
print(f"Salvato: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slide)")
